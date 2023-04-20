import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
import os

def extract_text(slide):
    text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += "%%" + run.text + "\n"
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text += "%%" + run.text + "\n"
        elif shape.has_chart:
            for series in shape.chart.series:
                text += "%%" + series.name + "\n"
                for point in series.points:
                    text += "%%" + point.name + "\n"
    return text


def save_text_to_file(text, file_name):
    with open(file_name, "w", encoding="utf-8") as file:
        file.write(text)

def process_pptx_file(file_path):
    prs = Presentation(file_path)
    extracted_text = "%%\n"
    
    for idx, slide in enumerate(prs.slides):
        extracted_text += f"[슬라이드-{idx + 1}]\n"
        extracted_text += extract_text(slide) + "\n"
    
    return extracted_text

def open_file():
    file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    if not file_path:
        return

    extracted_text = process_pptx_file(file_path)
    save_text_to_file(extracted_text, "output.txt")

app = tk.Tk()
app.title("PPTX Text Extractor")

frame = tk.Frame(app, padx=20, pady=20)
frame.pack()

button_open = tk.Button(frame, text="Open PowerPoint file", command=open_file)
button_open.pack()

app.mainloop()
